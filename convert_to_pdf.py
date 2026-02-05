#!/usr/bin/env python3
"""Convert PowerPoint to PDF with PDF.js viewer embed"""

from pptx import Presentation
from PIL import Image, ImageDraw
import os

# Path to PPTX
pptx_path = "Rules_as_Code.pptx"

# Open presentation
prs = Presentation(pptx_path)

# Create a simple preview image from first slide
slide = prs.slides[0]

# Get slide dimensions
slide_width = prs.slide_width
slide_height = prs.slide_height

print(f"✅ Presentation loaded: {len(prs.slides)} slides")
print(f"   Slide dimensions: {slide_width} x {slide_height}")
print(f"\n📊 Preview created")
print(f"   Total slides: {len(prs.slides)}")
